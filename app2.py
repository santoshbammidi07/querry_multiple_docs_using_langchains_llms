import streamlit as st
import base64
import json
import pandas as pd
from pathlib import Path
from mistralai import Mistral, DocumentURLChunk, ImageURLChunk, TextChunk
from langchain.vectorstores import FAISS
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.chains import ConversationalRetrievalChain
from langchain.chat_models import ChatOpenAI
from pdfminer.high_level import extract_text
from pptx import Presentation
from dotenv import load_dotenv

# Initialize Mistral client
api_key = "YOUR_MISTRAL_API_KEY"
client = Mistral(api_key=api_key)

def extract_pdf_text(pdf_file):
    return extract_text(pdf_file)

def extract_ppt_text(ppt_file):
    prs = Presentation(ppt_file)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text

def extract_excel_text(excel_file):
    df = pd.read_excel(excel_file, sheet_name=None)
    text = "\n".join(df[sheet].to_string() for sheet in df)
    return text

def extract_text_from_files(files):
    extracted_text = ""
    for file in files:
        file_extension = file.name.split(".")[-1].lower()
        if file_extension == "pdf":
            extracted_text += extract_pdf_text(file)
        elif file_extension == "pptx":
            extracted_text += extract_ppt_text(file)
        elif file_extension in ["xls", "xlsx"]:
            extracted_text += extract_excel_text(file)
    return extracted_text

def extract_images_from_pdf(pdf_file):
    uploaded_file = client.files.upload(
        file={"file_name": pdf_file.name, "content": pdf_file.getvalue()},
        purpose="ocr",
    )
    signed_url = client.files.get_signed_url(file_id=uploaded_file.id, expiry=1)
    response = client.ocr.process(
        document=DocumentURLChunk(document_url=signed_url.url), 
        model="mistral-ocr-latest", 
        include_image_base64=True
    )
    return response

def process_ocr_images(ocr_response):
    extracted_text = ""
    for page in ocr_response.pages:
        for img in page.images:
            extracted_text += img.markdown + "\n"
    return extracted_text

def get_vectorstore(text):
    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_texts([text], embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(model_name="gpt-4", temperature=0)
    return ConversationalRetrievalChain.from_llm(llm, vectorstore.as_retriever())

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']
    for i, message in enumerate(st.session_state.chat_history):
        st.write(f"{'User' if i % 2 == 0 else 'Bot'}: {message.content}")

def main():
    load_dotenv()
    st.set_page_config(page_title="Chat with Documents", page_icon=":books:")
    
    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None
    
    st.header("Chat with your Documents :books:")
    user_question = st.text_input("Ask a question about your documents:")
    if user_question:
        handle_userinput(user_question)
    
    with st.sidebar:
        st.subheader("Upload your documents")
        files = st.file_uploader("Upload PDF, PPTX, or Excel", accept_multiple_files=True)
        if st.button("Process"):
            with st.spinner("Processing"):
                raw_text = extract_text_from_files(files)
                
                for file in files:
                    if file.name.endswith("pdf"):
                        ocr_response = extract_images_from_pdf(file)
                        raw_text += process_ocr_images(ocr_response)
                
                vectorstore = get_vectorstore(raw_text)
                st.session_state.conversation = get_conversation_chain(vectorstore)

if __name__ == '__main__':
    main()
